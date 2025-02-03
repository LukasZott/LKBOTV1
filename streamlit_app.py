import streamlit as st
import openai
import os
import time
import json
import pandas as pd
import pdfplumber
import pptx
import openpyxl
from PIL import Image

# Configuração inicial da página
st.set_page_config(
    page_title="LBOT V1",
    page_icon="💛",
    layout="wide",
)

# CSS personalizado para estilizar o balão de upload e o aviso
st.markdown(
    """
    <style>
    /* Estilo para o texto na sidebar */
    .stSidebar .stMarkdown, .stSidebar .stTextInput, .stSidebar .stTextArea, .stSidebar .stButton, .stSidebar .stExpander {
        color: white !important;  /* Cor do texto na sidebar */
    }

    /* Estilo para o texto na parte principal */
    .stMarkdown, .stTextInput, .stTextArea, .stButton, .stExpander {
        color: black !important;  /* Cor do texto na parte principal */
    }

    /* Estilo para o container de upload de arquivos */
    .stFileUploader > div > div {
        background-color: white;  /* Fundo branco */
        color: black;  /* Texto preto */
        border-radius: 10px;
        padding: 10px;
        border: 1px solid #ccc;  /* Borda cinza para destacar */
    }

    /* Estilo para o texto dentro do balão de upload */
    .stFileUploader label {
        color: black !important;  /* Texto preto */
    }

    /* Estilo para o botão de upload */
    .stFileUploader button {
        background-color: #8dc50b;  /* Verde */
        color: white;  /* Texto branco */
        border-radius: 5px;
        border: none;
        padding: 8px 16px;
    }

    /* Estilo para o texto de drag and drop */
    .stFileUploader div[data-testid="stFileUploaderDropzone"] {
        color: black !important;  /* Texto preto */
    }

    /* Estilo para o container de avisos (st.warning) */
    div[data-testid="stNotification"] > div > div {
        background-color: white !important;  /* Fundo branco */
        color: black !important;  /* Texto preto */
        border-radius: 10px !important;
        padding: 10px !important;
        border: 1px solid #ccc !important;  /* Borda cinza para destacar */
    }

    /* Estilo para o ícone de aviso */
    div[data-testid="stNotification"] > div > div > div:first-child {
        color: #8dc50b !important;  /* Cor do ícone (verde) */
    }

    /* Estilo para o subtítulo */
    .subtitulo {
        font-size: 16px !important;  /* Tamanho da fonte reduzido */
        color: black !important;  /* Cor do texto alterada para preto */
    }

    /* Estilo para o rótulo do campo de entrada na sidebar */
    .stSidebar label {
        color: white !important;  /* Cor do texto branco */
    }

    /* Estilo para o texto na caixa de entrada do chat */
    .stChatInput input {
        color: white !important;  /* Cor do texto branco */
    }

    /* Estilo para o placeholder na caixa de entrada do chat */
    .stChatInput input::placeholder {
        color: white !important;  /* Cor do placeholder branco */
    }

    /* Estilo para o texto na caixa de entrada do chat */
div.stChatInput textarea {
    color: white !important;  /* Cor do texto branco */
}

/* Estilo para o placeholder na caixa de entrada do chat */
div.stChatInput textarea::placeholder {
    color: white !important;  /* Cor do placeholder branco */
    opacity: 1;  /* Garante que o placeholder seja totalmente visível */
}
    
     /* Estilo para o ícone */
    .stImage > img {
        filter: drop-shadow(0 0 0 #8dc50b);  /* Aplica a cor #8dc50b ao ícone */
    }
    </style>
    """,
    unsafe_allow_html=True
)


# Verificar se o arquivo do ícone existe
if os.path.exists(ICON_PATH):
    try:
        # Usar st.columns para posicionar o ícone ao lado do título
        col1, col2 = st.columns([1.5, 4])  # Ajuste as proporções conforme necessário
        with col1:
            st.image(ICON_PATH, width=10000000)  # Exibe o ícone com largura de 30px
        with col2:
            st.title("LKBOT")  # Exibe o título
    except Exception as e:
        st.error(f"Erro ao carregar o ícone: {e}")
else:
    st.title("LKBOT")  # Fallback se o ícone não existir

# Subtítulo com fonte reduzida e texto preto
st.markdown(
    '<cp class="subtitulo">Pronto para ajudar!.</p>',
    unsafe_allow_html=True
)

# Criar uma opção de seleção para armazenar arquivos
st.sidebar.subheader("📂 Configuração de Arquivos")

# Opções de local para armazenar os arquivos
opcoes_local = ["Pasta Padrão (uploads/)", "Escolher Diretório Personalizado"]
escolha_local = st.sidebar.selectbox("Escolha onde armazenar os arquivos:", opcoes_local)

# Campo de entrada para definir um diretório personalizado
diretorio_personalizado = None
if escolha_local == "Escolher Diretório Personalizado":
    diretorio_personalizado = st.sidebar.text_input("📁 Digite o caminho do diretório:")

# Diretório de upload
UPLOAD_FOLDER = "uploads"
if escolha_local == "Escolher Diretório Personalizado" and diretorio_personalizado:
    UPLOAD_FOLDER = diretorio_personalizado

# Criar diretório caso não exista
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# Função para salvar arquivos no diretório escolhido
def salvar_arquivo(arquivo):
    caminho_arquivo = os.path.join(UPLOAD_FOLDER, arquivo.name)
    with open(caminho_arquivo, "wb") as f:
        f.write(arquivo.getbuffer())
    return caminho_arquivo

# Interface para upload de arquivos
st.sidebar.subheader("📤 Upload de Documentos")
arquivos = st.sidebar.file_uploader(
    "Carregue arquivos (PDF, CSV, XLSX, PPTX)",
    type=["pdf", "csv", "xlsx", "pptx"],
    accept_multiple_files=True
)

# Processar e armazenar os arquivos carregados
documentos_carregados = []
if arquivos:
    for arquivo in arquivos:
        caminho = salvar_arquivo(arquivo)
        documentos_carregados.append(caminho)
    st.sidebar.success(f"Arquivos armazenados em: {UPLOAD_FOLDER}")

# Função para processar os arquivos armazenados
def processar_arquivos():
    contexto = ""
    
    for caminho in documentos_carregados:
        if caminho.endswith(".pdf"):
            with pdfplumber.open(caminho) as pdf:
                for page in pdf.pages:
                    contexto += page.extract_text() + "\n"
        elif caminho.endswith(".csv"):
            df = pd.read_csv(caminho)
            contexto += df.to_string() + "\n"
        elif caminho.endswith(".xlsx"):
            df = pd.read_excel(caminho)
            contexto += df.to_string() + "\n"
        elif caminho.endswith(".pptx"):
            prs = pptx.Presentation(caminho)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        contexto += shape.text + "\n"
    
    return contexto

# Processa os arquivos e adiciona ao contexto
contexto_documentos = processar_arquivos()

# Função para gerar resposta usando GPT-4o
def gerar_resposta(pergunta):
    if not contexto_documentos:
        return "Nenhum documento carregado para análise."

    contexto_pergunta = f"Baseado nos documentos carregados, responda: {pergunta}\n\n"
    contexto_pergunta += contexto_documentos[:2000]  # Limita o contexto para evitar excesso de tokens

    mensagens = [
        {"role": "system", "content": "Você é uma IA especializada em Administração Pública."},
        {"role": "user", "content": contexto_pergunta}
    ]

    tentativas = 3
    for tentativa in range(tentativas):
        try:
            time.sleep(1)
            resposta = openai.ChatCompletion.create(
                model="gpt-4o",
                messages=mensagens,
                temperature=0.3,
                max_tokens=800
            )
            return resposta["choices"][0]["message"]["content"]
        except Exception as e:
            if tentativa < tentativas - 1:
                time.sleep(2)
                continue
            else:
                return f"Erro ao gerar a resposta: {str(e)}"

# Entrada para perguntas no chat
user_input = st.chat_input("💬 Sua pergunta:")
if user_input and user_input.strip():
    resposta = gerar_resposta(user_input)
    st.write(f"**CADE IA:** {resposta}")
